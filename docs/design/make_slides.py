#
# Copyright (C) 2026 Advanced Micro Devices, Inc. All rights reserved.
# Licensed under the MIT License.
#
"""Generate the executive slide deck for the hybrid NPU + GPU project.

    python docs/design/make_slides.py

Writes docs/design/hybrid-npu-gpu-slides.pptx. Audience is leadership rather
than implementers: major implementation shape, no code, no per-task detail.
Content is derived from hybrid-npu-gpu-design.md and hybrid-npu-gpu-tasks.md;
update those first, then regenerate.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "hybrid-npu-gpu-slides.pptx"

FONT = "Segoe UI"

INK = "14181F"
BODY = "333B47"
MUTE = "6B7480"
LINE = "DDE2E8"
WHITE = "FFFFFF"

NPU = "B45309"
NPU_BG = "FFFBEB"
GPU = "1D4ED8"
GPU_BG = "EFF6FF"
OK = "047857"
OK_BG = "ECFDF5"
RED = "B91C1C"
RED_BG = "FEF4F4"
PANEL = "F8FAFC"

W, H = 13.333, 7.5


def rgb(hex_str):
    return RGBColor.from_string(hex_str)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def text(
    slide,
    x,
    y,
    w,
    h,
    paras,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", align)
        p.space_after = Pt(para.get("after", 6))
        p.space_before = Pt(para.get("before", 0))
        if "line" in para:
            p.line_spacing = para["line"]
        for run_spec in para["runs"]:
            body, size, bold, colour = (run_spec + (None,) * 4)[:4]
            run = p.add_run()
            run.text = body
            run.font.name = FONT
            run.font.size = Pt(size or 14)
            run.font.bold = bool(bold)
            run.font.color.rgb = rgb(colour or BODY)
    return tb


def box(slide, x, y, w, h, lines, fill=WHITE, border=None, lw=1.5, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if border:
        shape.line.color.rgb = rgb(border)
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.09)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        body, size, bold, colour = (line + (None,) * 4)[:4]
        run = p.add_run()
        run.text = body
        run.font.name = FONT
        run.font.size = Pt(size or 13)
        run.font.bold = bool(bold)
        run.font.color.rgb = rgb(colour or BODY)
    return shape


def arrow(slide, x, y, w, colour=MUTE, h=0.16):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colour)
    shape.line.fill.background()
    return shape


def rule(slide, x, y, w, colour=LINE, thickness=1.2):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(thickness)
    )
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(colour)
    shape.line.fill.background()
    return shape


def heading(slide, title, kicker=None):
    if kicker:
        text(
            slide,
            0.75,
            0.42,
            11.8,
            0.3,
            [{"runs": [(kicker.upper(), 11, True, MUTE)]}],
        )
        top = 0.72
    else:
        top = 0.55
    text(slide, 0.75, top, 11.8, 0.6, [{"runs": [(title, 30, True, INK)]}])
    rule(slide, 0.75, top + 0.72, 11.8)
    return top + 0.95


def footer(slide, page):
    text(
        slide,
        0.75,
        6.95,
        8.0,
        0.3,
        [{"runs": [("Hybrid NPU + GPU Execution  ·  MorphiZen HIP EP", 9, False, MUTE)]}],
    )
    text(
        slide,
        11.9,
        6.95,
        0.65,
        0.3,
        [{"runs": [(str(page), 9, False, MUTE)]}],
        align=PP_ALIGN.RIGHT,
    )


# --------------------------------------------------------------------- slides


def slide_title(prs):
    slide = add_slide(prs)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(2.45)
    )
    band.shadow.inherit = False
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(INK)
    band.line.fill.background()

    text(
        slide,
        0.9,
        0.85,
        11.5,
        0.8,
        [{"runs": [("Hybrid NPU + GPU Execution", 40, True, WHITE)]}],
    )
    text(
        slide,
        0.9,
        1.62,
        11.5,
        0.5,
        [
            {
                "runs": [
                    (
                        "Faster prompt response and lower power for on-device LLMs",
                        17,
                        False,
                        "C9D0D8",
                    )
                ]
            }
        ],
    )

    text(
        slide,
        0.9,
        3.0,
        11.5,
        0.9,
        [
            {
                "runs": [
                    ("Use each processor for the phase it is good at — ", 19, False, BODY),
                    ("the NPU for reading the prompt, the GPU for writing the answer", 19, True, INK),
                    (" — inside one unchanged application.", 19, False, BODY),
                ],
                "line": 1.35,
            }
        ],
    )

    for i, (label, value, colour) in enumerate(
        [
            ("Target", "Ryzen AI  ·  Strix-class APU", NPU),
            ("First models", "Llama-3.2-1B, then Llama-3.1-8B", GPU),
            ("Acceptance", "Accuracy, not a throughput number", OK),
        ]
    ):
        x = 0.9 + i * 3.95
        box(
            slide,
            x,
            4.45,
            3.6,
            1.0,
            [(label.upper(), 10, True, colour), (value, 12, False, BODY)],
            fill=PANEL,
            border=LINE,
            lw=1.0,
        )

    text(
        slide,
        0.9,
        6.3,
        11.5,
        0.4,
        [{"runs": [("Design review complete  ·  Phase 0 in progress", 12, False, MUTE)]}],
    )
    return slide


def slide_opportunity(prs):
    slide = add_slide(prs)
    top = heading(
        slide, "Two phases, opposite needs", kicker="The opportunity"
    )

    text(
        slide,
        0.75,
        top,
        11.8,
        0.5,
        [
            {
                "runs": [
                    ("Generating text has two stages that stress hardware in ", 15, False, BODY),
                    ("completely different ways", 15, True, INK),
                    (". Today we run both on the GPU, so one of them is always on the wrong processor.", 15, False, BODY),
                ],
                "line": 1.3,
            }
        ],
    )

    y = top + 0.85
    box(
        slide,
        0.75,
        y,
        5.7,
        2.05,
        [
            ("READING THE PROMPT", 11, True, NPU),
            ("Happens once. Processes every word at the same time.", 13, False, BODY),
            ("Limited by raw compute", 14, True, INK),
            ("Best on the NPU", 15, True, NPU),
        ],
        fill=NPU_BG,
        border=NPU,
        align=PP_ALIGN.LEFT,
    )
    box(
        slide,
        6.85,
        y,
        5.7,
        2.05,
        [
            ("WRITING THE ANSWER", 11, True, GPU),
            ("Repeats for every word produced. One word at a time.", 13, False, BODY),
            ("Limited by memory speed", 14, True, INK),
            ("Best on the GPU", 15, True, GPU),
        ],
        fill=GPU_BG,
        border=GPU,
        align=PP_ALIGN.LEFT,
    )

    box(
        slide,
        0.75,
        y + 2.35,
        11.8,
        0.95,
        [
            (
                "Splitting them improves how fast the first word appears, and frees each processor "
                "to idle while the other works.",
                15,
                True,
                INK,
            ),
        ],
        fill=OK_BG,
        border=OK,
        align=PP_ALIGN.LEFT,
    )
    footer(slide, 2)
    return slide


def slide_approach(prs):
    slide = add_slide(prs)
    top = heading(slide, "One session, two processors", kicker="The approach")

    text(
        slide,
        0.75,
        top,
        11.8,
        0.4,
        [
            {
                "runs": [
                    ("The application is unchanged. Work is routed per request inside our execution provider.", 15, False, BODY)
                ]
            }
        ],
    )

    y = top + 0.75
    box(slide, 0.75, y + 0.35, 1.5, 1.0, [("Prompt", 14, True, INK)], fill=PANEL, border=LINE, lw=1.0)
    arrow(slide, 2.4, y + 0.78, 0.5)

    box(
        slide,
        3.05,
        y,
        3.0,
        1.7,
        [
            ("NPU", 20, True, NPU),
            ("Reads the whole prompt", 12, False, BODY),
            ("Runs once per request", 12, False, BODY),
        ],
        fill=NPU_BG,
        border=NPU,
        lw=2.0,
    )
    arrow(slide, 6.2, y + 0.78, 0.5)

    box(
        slide,
        6.85,
        y + 0.15,
        2.3,
        1.4,
        [("Shared memory", 13, True, OK), ("handoff point", 11, False, OK)],
        fill=OK_BG,
        border=OK,
        lw=2.0,
    )
    arrow(slide, 9.3, y + 0.78, 0.5)

    box(
        slide,
        9.95,
        y,
        2.6,
        1.7,
        [
            ("GPU", 20, True, GPU),
            ("Writes the answer", 12, False, BODY),
            ("Loops per word", 12, False, BODY),
        ],
        fill=GPU_BG,
        border=GPU,
        lw=2.0,
    )

    text(
        slide,
        0.75,
        y + 2.0,
        11.8,
        0.4,
        [
            {
                "runs": [
                    ("The handoff is the shared memory itself. ", 15, True, INK),
                    ("Nothing is transferred, converted or reorganised when work moves between processors.", 15, False, BODY),
                ]
            }
        ],
    )

    for i, (title_, body_) in enumerate(
        [
            ("No application change", "Same model file, same interface, one session."),
            ("Decided per request", "A conversation can switch back and forth freely."),
            ("GPU path untouched", "Today's behaviour is preserved exactly."),
        ]
    ):
        x = 0.75 + i * 4.03
        box(
            slide,
            x,
            y + 2.55,
            3.75,
            0.95,
            [(title_, 13, True, INK), (body_, 11, False, MUTE)],
            fill=PANEL,
            border=LINE,
            lw=1.0,
            align=PP_ALIGN.LEFT,
        )
    footer(slide, 3)
    return slide


def slide_why_now(prs):
    slide = add_slide(prs)
    top = heading(
        slide, "Why this is possible on this hardware", kicker="The enabler"
    )

    text(
        slide,
        0.75,
        top,
        11.8,
        0.5,
        [
            {
                "runs": [
                    ("On these processors the NPU, GPU and CPU share the ", 15, False, BODY),
                    ("same physical memory", 15, True, INK),
                    (". That is what makes the split worth doing at all.", 15, False, BODY),
                ]
            }
        ],
    )

    y = top + 0.8
    box(
        slide,
        0.75,
        y,
        5.7,
        2.3,
        [
            ("ON A SEPARATE GRAPHICS CARD", 10, True, RED),
            ("Each processor owns its own memory.", 13, False, BODY),
            ("Handing work over means copying the data across.", 13, False, BODY),
            ("The copy costs more than the NPU saves.", 13, True, RED),
        ],
        fill=RED_BG,
        border=RED,
        align=PP_ALIGN.LEFT,
    )
    box(
        slide,
        6.85,
        y,
        5.7,
        2.3,
        [
            ("ON THIS SHARED-MEMORY DESIGN", 10, True, OK),
            ("All three processors address one pool of memory.", 13, False, BODY),
            ("Handing work over means agreeing where it is.", 13, False, BODY),
            ("Nothing moves, so the handoff is effectively free.", 13, True, OK),
        ],
        fill=OK_BG,
        border=OK,
        align=PP_ALIGN.LEFT,
    )

    box(
        slide,
        0.75,
        y + 2.6,
        11.8,
        1.0,
        [
            ("The catch we designed around", 12, True, INK),
            (
                "Shared memory does not mean every processor can reach every buffer automatically. "
                "Access is granted by how memory is allocated, so we verify provenance rather than trusting that a read happens to work.",
                12,
                False,
                BODY,
            ),
        ],
        fill=PANEL,
        border=LINE,
        lw=1.0,
        align=PP_ALIGN.LEFT,
    )
    footer(slide, 4)
    return slide


def slide_zero_copy(prs):
    slide = add_slide(prs)
    top = heading(slide, "One non-negotiable requirement", kicker="The constraint")

    box(
        slide,
        0.75,
        top,
        11.8,
        1.0,
        [
            (
                "No data is copied when execution moves between the NPU and the GPU.",
                20,
                True,
                INK,
            ),
        ],
        fill=RED_BG,
        border=RED,
        lw=2.0,
    )

    y = top + 1.3
    text(
        slide,
        0.75,
        y,
        11.8,
        0.4,
        [
            {
                "runs": [
                    ("Why it is a requirement and not a goal: ", 15, True, INK),
                    ("the shared data is the largest item in the system and grows with conversation length. Copying it would erase the benefit entirely.", 15, False, BODY),
                ]
            }
        ],
    )

    y += 0.75
    for i, (title_, body_) in enumerate(
        [
            (
                "The failure mode is silence",
                "A copy still produces the correct answer. It only shows up as lost performance, so testing accuracy can never find it.",
            ),
            (
                "So it is enforced, not reviewed",
                "Memory is accepted only from approved sources, and every test counts the copies and requires the count to be zero.",
            ),
            (
                "And it is a stop condition",
                "If no approach achieves it, we escalate and revisit the architecture rather than ship a version that copies.",
            ),
        ]
    ):
        x = 0.75 + i * 4.03
        box(
            slide,
            x,
            y,
            3.75,
            2.0,
            [(title_, 14, True, INK), (body_, 12, False, BODY)],
            fill=WHITE,
            border=LINE,
            lw=1.2,
            align=PP_ALIGN.LEFT,
        )
    footer(slide, 5)
    return slide


def slide_building(prs):
    slide = add_slide(prs)
    top = heading(slide, "What we are building", kicker="Implementation")

    text(
        slide,
        0.75,
        top,
        11.8,
        0.4,
        [
            {
                "runs": [
                    ("Six additions. The existing GPU compiler and runtime are extended, not replaced.", 15, False, BODY)
                ]
            }
        ],
    )

    items = [
        ("Compatibility layer", "A thin boundary to the NPU software stack, isolating its dependencies from ours.", NPU),
        ("Substitute for testing", "A stand-in for the NPU so most work is verified without hardware.", NPU),
        ("Shared memory pool", "Working space the NPU can reach, alongside the memory we already manage.", OK),
        ("Provenance registry", "Records where every buffer came from, so an invalid one is rejected outright.", OK),
        ("Compiler output", "Turns a claimed model into a sequence of NPU operations, or declines it.", GPU),
        ("Runtime executor", "Carries out that sequence, and falls back safely to the GPU at any point.", GPU),
    ]
    for i, (title_, body_, colour) in enumerate(items):
        col, row = i % 3, i // 3
        x = 0.75 + col * 4.03
        y = top + 0.7 + row * 1.65
        box(
            slide,
            x,
            y,
            3.75,
            1.45,
            [(title_, 14, True, colour), (body_, 11.5, False, BODY)],
            fill=WHITE,
            border=LINE,
            lw=1.2,
            align=PP_ALIGN.LEFT,
        )

    box(
        slide,
        0.75,
        top + 4.05,
        11.8,
        0.85,
        [
            (
                "Deliberate choice: most of this is verifiable on an ordinary developer machine. Only correctness and performance need the NPU.",
                14,
                True,
                INK,
            )
        ],
        fill=PANEL,
        border=LINE,
        lw=1.0,
        align=PP_ALIGN.LEFT,
    )
    footer(slide, 6)
    return slide


def slide_plan(prs):
    slide = add_slide(prs)
    top = heading(slide, "Delivery plan", kicker="Sequencing")

    text(
        slide,
        0.75,
        top,
        11.8,
        0.4,
        [
            {
                "runs": [
                    ("Nine stages. Every stage carries at least one formal checkpoint that must be signed off before the next begins.", 15, False, BODY)
                ]
            }
        ],
    )

    # Gate counts come from the work graph; keep them in step with it.
    stages = [
        ("0", "Prove the risky assumptions", 2),
        ("1", "Boundary to the NPU stack", 1),
        ("2", "Shared memory and enforcement", 2),
        ("3", "Compiler produces a plan", 1),
        ("4", "Routing and execution", 2),
        ("5", "Operation coverage", 1),
        ("6", "Handling unsupported operations", 1),
        ("7", "Loading model weights", 1),
        ("8", "Full end-to-end generation", 1),
    ]
    y = top + 0.62
    for i, (num, label, gates) in enumerate(stages):
        yy = y + i * 0.47
        box(
            slide,
            0.75,
            yy,
            0.58,
            0.4,
            [(num, 14, True, WHITE)],
            fill=INK,
        )
        text(
            slide,
            1.5,
            yy + 0.06,
            6.6,
            0.32,
            [{"runs": [(label, 14, False, BODY)]}],
        )
        box(
            slide,
            8.2,
            yy + 0.02,
            1.9,
            0.36,
            [(f"{gates} checkpoint" + ("s" if gates > 1 else ""), 10, True, NPU)],
            fill=NPU_BG,
            border=NPU,
            lw=1.0,
        )

    for i, (value, label, colour, bg) in enumerate(
        [
            ("39", "tasks in total", INK, PANEL),
            ("12", "checkpoints, any of which can halt the project", NPU, NPU_BG),
            ("20", "of those tasks need the NPU machine to sign off", GPU, GPU_BG),
        ]
    ):
        box(
            slide,
            10.45,
            y + i * 1.42,
            2.1,
            1.3,
            [(value, 26, True, colour), (label, 10, False, MUTE)],
            fill=bg,
            border=colour if colour != INK else LINE,
            lw=1.0,
        )
    footer(slide, 7)
    return slide


def slide_status(prs):
    slide = add_slide(prs)
    top = heading(slide, "Where we are today", kicker="Status")

    for i, (state, label, detail, colour, bg) in enumerate(
        [
            (
                "Complete",
                "Design and task breakdown",
                "Reviewed. Memory layouts confirmed compatible, which retired the largest technical risk.",
                OK,
                OK_BG,
            ),
            (
                "Awaiting hardware",
                "Two Phase 0 checks",
                "Software is written and reviewed. Both need the NPU machine to run, which is a scheduled trip.",
                NPU,
                NPU_BG,
            ),
            (
                "Blocked",
                "Everything after Phase 0",
                "Held deliberately. The memory question decides the architecture, so we do not build ahead of it.",
                MUTE,
                PANEL,
            ),
        ]
    ):
        x = 0.75 + i * 4.03
        box(
            slide,
            x,
            top,
            3.75,
            2.25,
            [
                (state.upper(), 10, True, colour),
                (label, 15, True, INK),
                (detail, 12, False, BODY),
            ],
            fill=bg,
            border=colour if colour != MUTE else LINE,
            lw=1.5,
            align=PP_ALIGN.LEFT,
        )

    box(
        slide,
        0.75,
        top + 2.6,
        11.8,
        1.55,
        [
            ("The one thing that would change the plan", 15, True, INK),
            (
                "Zero copy depends on the NPU being able to use memory our software allocated. "
                "That has never been tested in this combination. It is the first thing we check, on purpose — "
                "a negative result is an architecture decision, not a delay, and we would rather find it now than in Phase 5.",
                13,
                False,
                BODY,
            ),
        ],
        fill=WHITE,
        border=RED,
        lw=2.0,
        align=PP_ALIGN.LEFT,
    )
    footer(slide, 8)
    return slide


def slide_risks(prs):
    slide = add_slide(prs)
    top = heading(slide, "Risks, and how we find out early", kicker="Risk")

    text(
        slide,
        0.75,
        top,
        11.8,
        0.4,
        [
            {
                "runs": [
                    ("Three findings would invalidate the approach rather than complicate it. All three are checked before the bulk of the work.", 15, False, BODY)
                ]
            }
        ],
    )

    y = top + 0.62
    for i, (risk, detect, when) in enumerate(
        [
            (
                "Shared memory cannot be shared",
                "The NPU rejects memory our software allocated, so a copy would be unavoidable.",
                "Checked in Phase 0",
            ),
            (
                "The software stacks conflict",
                "Our runtime and the NPU's cannot coexist in one process. We have been bitten by this class of problem before.",
                "Checked in Phase 0",
            ),
            (
                "Too much falls back to the GPU",
                "If the NPU cannot handle enough of the model, the arrangement loses its benefit.",
                "Detected automatically at build time",
            ),
        ]
    ):
        yy = y + i * 1.18
        box(
            slide,
            0.75,
            yy,
            4.4,
            1.05,
            [(risk, 14, True, INK)],
            fill=RED_BG,
            border=RED,
            lw=1.2,
            align=PP_ALIGN.LEFT,
        )
        text(
            slide,
            5.45,
            yy + 0.1,
            4.5,
            0.85,
            [{"runs": [(detect, 12, False, BODY)], "line": 1.22}],
        )
        box(
            slide,
            10.2,
            yy + 0.25,
            2.35,
            0.55,
            [(when, 11, True, OK)],
            fill=OK_BG,
            border=OK,
            lw=1.0,
        )

    box(
        slide,
        0.75,
        y + 3.72,
        11.8,
        0.78,
        [
            (
                "Each checkpoint states who can sign it off. Work that needs the NPU cannot be closed by a successful run on a developer machine.",
                14,
                True,
                INK,
            )
        ],
        fill=PANEL,
        border=LINE,
        lw=1.0,
        align=PP_ALIGN.LEFT,
    )
    footer(slide, 9)
    return slide


def slide_success(prs):
    slide = add_slide(prs)
    top = heading(slide, "What success looks like", kicker="Definition of done")

    for i, (headline, detail) in enumerate(
        [
            (
                "Both models run split across NPU and GPU",
                "Llama-3.2-1B and Llama-3.1-8B, in one session, through one execution provider.",
            ),
            (
                "Nothing is copied, and we can prove it",
                "Copy counters read zero across a full generation, with no debug path active.",
            ),
            (
                "Answers match a trusted reference",
                "Output matches a CPU reference run; quality measured within tolerance.",
            ),
            (
                "Today's performance is not regressed",
                "Word-generation speed and accuracy are unchanged from the current GPU-only path.",
            ),
        ]
    ):
        col, row = i % 2, i // 2
        x = 0.75 + col * 6.1
        y = top + row * 1.72
        box(
            slide,
            x,
            y,
            5.7,
            1.5,
            [(headline, 15, True, INK), (detail, 12, False, BODY)],
            fill=WHITE,
            border=OK,
            lw=1.5,
            align=PP_ALIGN.LEFT,
        )

    box(
        slide,
        0.75,
        top + 3.6,
        11.8,
        1.15,
        [
            ("Accuracy is the acceptance criterion", 16, True, INK),
            (
                "No stage is signed off on a speed number. Performance work follows once the arrangement is correct — "
                "which keeps us from optimising something that is subtly wrong.",
                13,
                False,
                BODY,
            ),
        ],
        fill=OK_BG,
        border=OK,
        lw=1.5,
        align=PP_ALIGN.LEFT,
    )
    footer(slide, 10)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    slide_title(prs)
    slide_opportunity(prs)
    slide_approach(prs)
    slide_why_now(prs)
    slide_zero_copy(prs)
    slide_building(prs)
    slide_plan(prs)
    slide_status(prs)
    slide_risks(prs)
    slide_success(prs)

    prs.save(OUT)
    print(f"  wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
